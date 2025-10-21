import sys
import os
import subprocess
from PyQt6.QtWidgets import (
    QApplication, QWidget, QPushButton, QLabel, QFileDialog, QVBoxLayout,
    QHBoxLayout, QGridLayout, QScrollArea, QProgressBar, QMessageBox, QSpinBox, 
    QSizePolicy, QProgressDialog, QComboBox, QInputDialog, QTabWidget
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
from PyQt6.QtGui import QPixmap, QIcon, QPalette, QColor
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import license_manager

class ImageLoaderThread(QThread):
    progress_changed = pyqtSignal(int, int, str)  # current, total, tag
    finished_loading = pyqtSignal(list, dict, dict, dict)  # samples, tag_map, loaded_images_pixmap, loaded_images_pil

    def __init__(self, base_path, max_width=350):
        super().__init__()
        self.base_path = base_path
        self.max_width = max_width

    def run(self):
        samples = []
        tag_map = {}

        ed_sample_paths = []
        for root, dirs, _ in os.walk(self.base_path):
            for d in dirs:
                if d.upper().startswith("ED"):
                    ed_sample_paths.append(os.path.join(root, d))

        samples = [os.path.basename(p) for p in ed_sample_paths]

        for sample_path in ed_sample_paths:
            sample = os.path.basename(sample_path)
            for root, _, files in os.walk(sample_path):
                for file in files:
                    if file.lower().endswith(('.jpg', '.jpeg', '.png')):
                        tag = file.lower()
                        tag_map.setdefault(tag, {})[sample] = os.path.join(root, file)

        tags = sorted(tag_map.keys())
        total_tags = len(tags)

        loaded_images_pixmap = {}
        loaded_images_pil = {}

        for i, tag in enumerate(tags, 1):
            loaded_images_pixmap[tag] = {}
            loaded_images_pil[tag] = {}
            for sample in samples:
                path = tag_map[tag].get(sample)
                if path:
                    pil_img = self.load_pil_image(path)
                    pixmap = self.pil_to_pixmap(pil_img) if pil_img else None
                    loaded_images_pixmap[tag][sample] = (pixmap, path)
                    loaded_images_pil[tag][sample] = pil_img
                else:
                    loaded_images_pixmap[tag][sample] = (None, None)
                    loaded_images_pil[tag][sample] = None
            self.progress_changed.emit(i, total_tags, tag)

        self.finished_loading.emit(samples, tag_map, loaded_images_pixmap, loaded_images_pil)

    def load_pil_image(self, path):
        try:
            img = Image.open(path)
            wpercent = (self.max_width / float(img.size[0]))
            height = int((float(img.size[1]) * float(wpercent)))
            img = img.resize((self.max_width, height), Image.LANCZOS)
            return img
        except Exception as e:
            print(f"Error loading image {path}: {e}")
            return None

    def pil_to_pixmap(self, pil_img):
        if pil_img is None:
            return None
        bio = BytesIO()
        pil_img.save(bio, format='PNG')
        bio.seek(0)
        pixmap = QPixmap()
        pixmap.loadFromData(bio.read())
        return pixmap

class ClickableLabel(QLabel):
    clicked = pyqtSignal()
    def mousePressEvent(self, event):
        self.clicked.emit()
        super().mousePressEvent(event)

class InspectoApp(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Inspecto Tool - visual comparison")

        # --- Window sizing ---
        screen = QApplication.primaryScreen()
        available_size = screen.availableGeometry().size()
        self.resize(int(available_size.width() * 1), int(available_size.height() * 0.95))

        # --- Main layout ---
        self.main_layout = QVBoxLayout(self)

        # --- Tabs ---
        self.tabs = QTabWidget()
        self.main_layout.addWidget(self.tabs)

        # --- Controls layout for folder tab ---
        self.controls_layout = QHBoxLayout()
        self.folder_label = QLabel("Select folder:")
        self.select_button = QPushButton("Select folder")
        self.load_button = QPushButton("Load")
        self.clear_button = QPushButton("Clear")
        self.clear_button.setEnabled(False)

        self.max_columns_label = QLabel("Max columns:")
        self.max_columns_label.setFixedWidth(75)
        self.max_columns_spin = QSpinBox()
        self.max_columns_spin.setRange(1, 10)
        self.max_columns_spin.setValue(4)
        self.max_columns_spin.setFixedWidth(50)

        self.img_width_label = QLabel("Image width:")
        self.img_width_label.setFixedWidth(80)
        self.img_width_spin = QSpinBox()
        self.img_width_spin.setRange(50, 1000)
        self.img_width_spin.setValue(350)
        self.img_width_spin.setFixedWidth(60)

        self.export_pdf_button = QPushButton("Export to PowerPoint")
        if not license_manager.is_pro():
            self.export_pdf_button.setToolTip("Pro feature – activate license to enable Export")

        # Tag selection and jump
        self.tag_combo = QComboBox()
        self.tag_combo.setFixedWidth(200)
        self.jump_button = QPushButton("Skip to tag")
        self.jump_button.setEnabled(False)
        self.jump_button.clicked.connect(self.scroll_to_tag)

        # Add widgets to controls layout
        self.controls_layout.addWidget(self.folder_label)
        self.controls_layout.addWidget(self.select_button)
        self.controls_layout.addWidget(self.load_button)
        self.controls_layout.addWidget(self.clear_button)
        self.controls_layout.addStretch()
        self.controls_layout.addWidget(self.max_columns_label)
        self.controls_layout.addWidget(self.max_columns_spin)
        self.controls_layout.addWidget(self.img_width_label)
        self.controls_layout.addWidget(self.img_width_spin)
        self.controls_layout.addWidget(self.tag_combo)
        self.controls_layout.addWidget(self.jump_button)
        self.controls_layout.addWidget(self.export_pdf_button)

        self.activate_button = QPushButton("Activate Pro")
        self.activate_button.clicked.connect(self.activate_pro)
        self.controls_layout.addWidget(self.activate_button)

        # Subscription status
        self.pro_status_label = QLabel("")
        self.controls_layout.addWidget(self.pro_status_label)
        self.update_pro_status()

        # --- Folder tab ---
        self.folder_tab = QWidget()
        self.folder_tab_layout = QVBoxLayout(self.folder_tab)
        self.folder_tab_layout.addLayout(self.controls_layout)

        # Progress bar and status
        self.progress_bar = QProgressBar()
        self.progress_bar.setValue(0)
        self.progress_bar.hide()
        self.status_label = QLabel("")
        self.folder_tab_layout.addWidget(self.progress_bar)
        self.folder_tab_layout.addWidget(self.status_label)

        # Scroll area for images
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.images_widget = QWidget()
        self.grid_layout = QGridLayout(self.images_widget)
        self.grid_layout.setContentsMargins(5, 5, 5, 5)
        self.grid_layout.setSpacing(10)
        self.scroll_area.setWidget(self.images_widget)
        self.folder_tab_layout.addWidget(self.scroll_area)

        # Scrollbar styling
        scroll_style = """
        QScrollBar::handle:vertical, QScrollBar::handle:horizontal {
            background: #2a82da !important;
            border-radius: 4px;
            min-height: 20px;
            min-width: 20px;
        }
        QScrollBar:vertical, QScrollBar:horizontal {
            background: #353535;
        }
        """
        self.scroll_area.setStyleSheet(scroll_style)

        # Add folder tab to tabs
        self.tabs.addTab(self.folder_tab, "Folder View")

        # --- Custom Images tab ---
        self.custom_tab = QWidget()
        self.custom_tab_layout = QVBoxLayout(self.custom_tab)

        from custom_tab import CustomImageGrid
        self.custom_grid = CustomImageGrid()
        self.custom_grid.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)

        # --- Wrap grid in scroll area with dynamic stretch ---
        self.custom_scroll = QScrollArea()
        self.custom_scroll.setWidgetResizable(True)  # allows horizontal stretch
        self.custom_scroll.setWidget(self.custom_grid)
        self.custom_tab_layout.addWidget(self.custom_scroll)

        # Status label for custom tab
        self.custom_status_label = QLabel("Drag images here")
        self.custom_tab_layout.addWidget(self.custom_status_label)

        # Connect signal for dropped images
        self.custom_grid.images_dropped.connect(self.on_custom_images_loaded)

        # Add custom tab to tabs
        self.tabs.addTab(self.custom_tab, "Custom Images")

        # --- Internal state ---
        self.selected_folder = None
        self.tag_widgets = {}
        self.image_loader_thread = None
        self.loaded_images_pil_cache = {}

        # --- Signals for folder tab ---
        self.select_button.clicked.connect(self.select_folder)
        self.load_button.clicked.connect(self.load_images)
        self.clear_button.clicked.connect(self.clear_images)
        self.export_pdf_button.clicked.connect(self.on_export_clicked)





    def update_pro_status(self):
        if license_manager.is_pro():
            self.pro_status_label.setText("Status: Pro ✅")
            self.pro_status_label.setStyleSheet("color: green; font-weight: bold;")
        else:
            self.pro_status_label.setText("Status: Free 🔒")
            self.pro_status_label.setStyleSheet("color: red; font-weight: bold;")

    def unlock_pro_features(self):
        self.export_pdf_button.setEnabled(True)
        self.export_pdf_button.setToolTip("")  # clear tooltip
        self.update_pro_status()
        QMessageBox.information(self, "Pro Unlocked", "All Pro features are now available.")


    def activate_pro(self):
        if license_manager.is_pro():
            QMessageBox.information(self, "Already Pro",
                "You already have a valid Pro license! 🎉")
            return

        key, ok = QInputDialog.getText(self, "Activate Inspecto Pro", "Enter your license key:")
        if ok and key:
            key = key.strip().upper()
            if not license_manager.validate_key_format(key):
                QMessageBox.warning(self, "Invalid Key", "The license key format is incorrect.")
                return

            if not license_manager.verify_key_offline(key):
                QMessageBox.warning(self, "Invalid Key", "This license key is not valid.")
                return

            if not license_manager.verify_hwid_match(key):
                QMessageBox.warning(self, "Wrong Machine",
                    "This key is already used on another machine.\n"
                    "Each license is tied to a single computer.")
                return

            # Save license and show confirmation with HWID
            license_manager.save_license(key)
            QMessageBox.information(self, "Activated",
                f"Inspecto Pro activated! Saved to this machine.\nHWID: {license_manager.machine_fingerprint()}")
            self.unlock_pro_features()


    def on_export_clicked(self):
        if not license_manager.is_pro():
            QMessageBox.warning(self, "Pro Feature", "Export to PowerPoint is a Pro feature. Please activate your license.")
            return

        if not self.loaded_images_pil_cache:
            QMessageBox.warning(self, "No images", "Please load images before exporting.")
            return

        self.export_to_pptx()


    def select_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "Select folder", "")
        if folder:
            self.selected_folder = folder
            self.folder_label.setText(f"Folder: {folder}    ")

    def load_images(self):
        self.setFocus()
        if not self.selected_folder or not os.path.isdir(self.selected_folder):
            QMessageBox.warning(self, "Failure", "Incorrect path to folder.")
            return

        self.progress_bar.setValue(0)
        self.progress_bar.show()
        self.status_label.show()
        self.status_label.setText("Reading data and images...")
        self.load_button.setEnabled(False)
        self.clear_button.setEnabled(False)
        self.export_pdf_button.setEnabled(False)

        max_width = self.img_width_spin.value()
        self.image_loader_thread = ImageLoaderThread(self.selected_folder, max_width)
        self.image_loader_thread.progress_changed.connect(self.on_progress_changed)
        self.image_loader_thread.finished_loading.connect(self.on_finished_loading)
        self.image_loader_thread.start()

    def on_progress_changed(self, current, total, tag):
        percent = int(current / total * 100)
        self.progress_bar.setValue(percent)
        self.status_label.setText(f"Reading tag: {tag} ({current}/{total})")

    def on_finished_loading(self, samples, tag_map, loaded_images_pixmap, loaded_images_pil):
        self.status_label.setText("Reading done, building view...")
        self.progress_bar.setValue(0)
        self.progress_bar.hide()
        self.status_label.hide()

        while self.grid_layout.count():
            item = self.grid_layout.takeAt(0)
            widget = item.widget()
            if widget:
                widget.deleteLater()

        self.tag_widgets.clear()
        self.tag_combo.clear()

        max_columns = self.max_columns_spin.value()

        row = 0
        for idx, tag in enumerate(sorted(tag_map.keys())):
            tag_container = QWidget()
            tag_layout = QGridLayout(tag_container)
            tag_layout.setContentsMargins(10,10,10,10)
            tag_layout.setSpacing(10)

            tag_container.setStyleSheet("background-color: #ADD8E6;")

            tag_label = QLabel(f"tag: {tag}")
            tag_label.setAlignment(Qt.AlignmentFlag.AlignLeft | Qt.AlignmentFlag.AlignVCenter)
            tag_label.setStyleSheet("""
                font-weight: bold;
                color: black;
                font-size: 20px;
                background-color: #e0eaff;
                padding: 5px;
            """)
            tag_layout.addWidget(tag_label, 0, 0, 1, max_columns)

            col = 0
            tag_row = 1
            for sample in samples:
                pixmap, path = loaded_images_pixmap[tag].get(sample, (None, None))

                container = QWidget()
                container_layout = QVBoxLayout(container)
                container_layout.setContentsMargins(0,0,0,0)
                container_layout.setSpacing(0)
                container_layout.setAlignment(Qt.AlignmentFlag.AlignHCenter)
                container.setStyleSheet("border: 1px solid black;")

                if pixmap:
                    img_label = ClickableLabel()
                    img_label.setPixmap(pixmap)
                    img_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
                    img_label.setFixedWidth(pixmap.width())
                    img_label.clicked.connect(lambda checked=False, p=path: self.open_external_viewer(p))
                else:
                    placeholder_width = self.img_width_spin.value()
                    placeholder_height = 100
                    img_label = QLabel("No image")
                    img_label.setFixedSize(placeholder_width, placeholder_height)
                    img_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

                img_label.setStyleSheet("""
                    QLabel {
                        border: 2px solid black;
                        border-radius: 5px;
                        transition: all 0.2s ease;
                    }
                    QLabel:hover {
                        border: 2px solid red;
                        background-color: rgba(42, 130, 218, 0.1);
                    }
                """)

                sample_label = QLabel(sample)
                sample_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
                sample_label.setStyleSheet("color: black; font-size: 18px; border: 2px solid black; background-color: white;")
                sample_label.setSizePolicy(QSizePolicy.Policy.Preferred, QSizePolicy.Policy.Fixed)

                container_layout.addWidget(img_label)
                container_layout.addWidget(sample_label)

                tag_layout.addWidget(container, tag_row, col)

                col += 1
                if col >= max_columns:
                    col = 0
                    tag_row += 1

            self.grid_layout.addWidget(tag_container, row, 0)
            self.tag_widgets[tag] = tag_container
            self.tag_combo.addItem(tag)
            row += 1

        if self.tag_widgets:
            self.jump_button.setEnabled(True)

        self.clear_button.setEnabled(True)
        self.load_button.setEnabled(True)
        self.export_pdf_button.setEnabled(True)
        self.loaded_images_pil_cache = loaded_images_pil

    def clear_images(self):
        self.setFocus()
        while self.grid_layout.count():
            item = self.grid_layout.takeAt(0)
            widget = item.widget()
            if widget:
                widget.deleteLater()
        self.status_label.setText("")
        self.progress_bar.setValue(0)
        self.clear_button.setEnabled(False)
        self.export_pdf_button.setEnabled(False)
        self.jump_button.setEnabled(False)
        self.tag_combo.clear()
        self.tag_widgets.clear()

    def open_external_viewer(self, filepath):
        if not filepath or not os.path.isfile(filepath):
            QMessageBox.warning(self, "Failure", "File doesn´t exist or is incorrect.")
            return

        normalized_path = os.path.normpath(filepath)

        try:
            if sys.platform.startswith('win'):
                os.startfile(normalized_path)
            elif sys.platform.startswith('darwin'):
                subprocess.run(['open', normalized_path])
            else:
                subprocess.run(['xdg-open', normalized_path])
        except Exception as e:
            QMessageBox.warning(self, "Failure", f"Unable to open file: {e}")

    def scroll_to_tag(self):
        selected_tag = self.tag_combo.currentText()
        if not selected_tag or selected_tag not in self.tag_widgets:
            return
        widget = self.tag_widgets[selected_tag]
        self.scroll_area.ensureWidgetVisible(widget)

    def export_to_pptx(self):
        if not self.loaded_images_pil_cache:
            QMessageBox.warning(self, "Failure", "No images to export.")
            return

        filename, _ = QFileDialog.getSaveFileName(self, "Save PowerPoint", "", "PowerPoint Files (*.pptx)")
        if not filename:
            return

        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        margin = Inches(0.5)
        padding = Inches(0.2)
        header_height = Inches(0.7)

        usable_width = slide_width - 2 * margin
        usable_height = slide_height - 2 * margin - header_height

        max_columns = self.max_columns_spin.value()
        img_width = (usable_width - (max_columns - 1) * padding) / max_columns

        tags = sorted(self.loaded_images_pil_cache.keys())

        progress = QProgressDialog("Exporting to PowerPoint...", "Cancel", 0, len(tags), self)
        progress.setWindowTitle("Export PowerPoint")
        progress.setWindowModality(Qt.WindowModality.ApplicationModal)
        progress.setMinimumDuration(0)

        for i, tag in enumerate(tags):
            if progress.wasCanceled():
                break

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # prázdný slide


            # Tag vlevo nahoře
            header_height = Inches(0.4)
            tag_box = slide.shapes.add_textbox(margin, margin, usable_width, header_height)

            # Nastavení výplně na světle šedou (třeba RGB 220,220,220)
            fill = tag_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(220, 220, 220)

            # Nastavení černého rámečku
            line = tag_box.line
            line.color.rgb = RGBColor(0, 0, 0)
            line.width = Pt(1)  # tloušťka rámečku

            tf2 = tag_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = 1  # zarovnání vlevo
            run2 = p2.add_run()
            run2.text = f"tag: {tag}"
            font2 = run2.font
            font2.size = Pt(16)
            font2.bold = False
            font2.color.rgb = RGBColor(0, 0, 0)


            samples = list(self.loaded_images_pil_cache[tag].keys())

            rows = (len(samples) + max_columns - 1) // max_columns
            img_height = (usable_height - (rows - 1) * padding) / rows

            col = 0
            row = 0
            
            extra_top_padding = Inches(0.3)  # pevná mezera pod tagem

            for sample in samples:
                pil_img = self.loaded_images_pil_cache[tag][sample]
                if pil_img is None:
                    continue

                w, h = pil_img.size
                ratio = w / h
                draw_width = img_width
                draw_height = img_width / ratio

                if draw_height > img_height:
                    draw_height = img_height
                    draw_width = img_height * ratio

                x = margin + col * (img_width + padding)
                y = margin + header_height + extra_top_padding + row * (img_height + padding)

                img_stream = BytesIO()
                pil_img.save(img_stream, format='PNG')
                img_stream.seek(0)

                slide.shapes.add_picture(img_stream, x, y, width=draw_width, height=draw_height)

                # Label vystředěný přesně pod obrázkem
                text_box = slide.shapes.add_textbox(x, y + draw_height, draw_width, Inches(0.3))
                tf_sample = text_box.text_frame
                tf_sample.margin_left = 0
                tf_sample.margin_right = 0
                tf_sample.margin_top = 0
                tf_sample.margin_bottom = 0

                for p in tf_sample.paragraphs:
                    p.alignment = 1  # CENTER

                p_sample = tf_sample.paragraphs[0]
                run_sample = p_sample.add_run()
                run_sample.text = sample
                font_sample = run_sample.font
                font_sample.size = Pt(10)
                font_sample.color.rgb = RGBColor(0, 0, 0)

                col += 1
                if col >= max_columns:
                    col = 0
                    row += 1

            progress.setValue(i + 1)
        try:
            prs.save(filename)
        except Exception as e:
            QMessageBox.warning(self, "Failure", f"Unable to save PowerPoint: {e}")
            return

        try:
            if sys.platform.startswith('win'):
                os.startfile(filename)
            elif sys.platform.startswith('darwin'):
                subprocess.run(['open', filename])
            else:
                subprocess.run(['xdg-open', filename])
        except Exception as e:
            print(f"Unable to open PowerPoint: {e}")

        QMessageBox.information(self, "Done", f"PowerPoint slides saved to:\n{filename}")

    def on_custom_images_loaded(self, image_paths):
        """Handle images dropped in the Custom Images tab."""
        print("Dropped images:", image_paths)
        # Optionally, update something like a counter or status label:
        self.status_label.setText(f"{len(image_paths)} custom images loaded")


if __name__ == "__main__":
    
    if hasattr(sys, '_MEIPASS'):
        icon_path = os.path.join(sys._MEIPASS, 'app_icon.ico')
    else:
        icon_path = 'app_icon.ico'

    app = QApplication(sys.argv)
    app.setWindowIcon(QIcon(icon_path))

    app.setStyle("Fusion")
    dark_palette = QPalette()
    dark_palette.setColor(QPalette.ColorRole.Window, QColor(53, 53, 53))
    dark_palette.setColor(QPalette.ColorRole.WindowText, QColor(220, 220, 220))
    dark_palette.setColor(QPalette.ColorRole.Base, QColor(35, 35, 35))
    dark_palette.setColor(QPalette.ColorRole.AlternateBase, QColor(53, 53, 53))
    dark_palette.setColor(QPalette.ColorRole.ToolTipBase, QColor(255, 255, 255))
    dark_palette.setColor(QPalette.ColorRole.ToolTipText, QColor(0, 0, 0))
    dark_palette.setColor(QPalette.ColorRole.Text, QColor(220, 220, 220))
    dark_palette.setColor(QPalette.ColorRole.Button, QColor(53, 53, 53))
    dark_palette.setColor(QPalette.ColorRole.ButtonText, QColor(220, 220, 220))
    dark_palette.setColor(QPalette.ColorRole.BrightText, QColor(255, 0, 0))
    dark_palette.setColor(QPalette.ColorRole.Highlight, QColor(42, 130, 218))
    dark_palette.setColor(QPalette.ColorRole.HighlightedText, QColor(255, 255, 255))

    app.setPalette(dark_palette)

    window = InspectoApp()
    window.show()
    sys.exit(app.exec())